from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Student Dropout Predictor"
subtitle.text = "Overview & Usage — Auto-generated"

# Slide 2: Project Overview
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.shapes.placeholders[1].text_frame

title.text = "Project Overview"
body.text = "Goal: Predict student dropout probability using logistic regression and provide an interactive Streamlit dashboard."

p = body.add_paragraph()
p.text = "Dataset: student dropout data (demographics, grades, behavior, support, etc.)"
p.level = 1

# Slide 3: Data & Features
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.shapes.placeholders[1].text_frame

title.text = "Key Data & Features"
body.text = "Examples of features included:"
for f in [
    'School, Gender, Age, Address',
    'Family_Size, Parental_Status',
    'Mother_Education, Father_Education',
    'Study_Time, Number_of_Failures, Absences',
    'Grades (Grade_1, Grade_2, Final_Grade)',
    'Support, Internet_Access, Extra Activities'
]:
    p = body.add_paragraph()
    p.text = f
    p.level = 1

# Slide 4: Data Processing
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.shapes.placeholders[1].text_frame

title.text = "Data Preprocessing"
body.text = "Main steps:"
steps = [
    'Load CSV and drop missing values',
    "Map 'yes'/'no' to 1/0 and encode categorical fields",
    'Binary encode School, Gender, Address, Parental_Status, Family_Size',
    'One-hot encode job, reason, guardian columns',
    'Return processed DataFrame for training and app'
]
for s in steps:
    p = body.add_paragraph()
    p.text = s
    p.level = 1

# Slide 5: Model
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.shapes.placeholders[1].text_frame

title.text = "Model"
body.text = "Logistic Regression model (scikit-learn)"
items = [
    'Pipeline: StandardScaler + LogisticRegression',
    'max_iter increased to 1000 to ensure convergence',
    'Model saved: model/logistic_regression_model.pkl'
]
for it in items:
    p = body.add_paragraph()
    p.text = it
    p.level = 1

# Slide 6: Streamlit App
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.shapes.placeholders[1].text_frame

title.text = "Streamlit Dashboard"
body.text = "Features:"
for line in [
    'Interactive sidebar inputs for student attributes',
    'Predicts probability (Not Dropout / Dropout) and final message',
    'Shows model performance: Confusion Matrix and ROC curve',
    'Correlation plots and heatmap'
]:
    p = body.add_paragraph()
    p.text = line
    p.level = 1

# Slide 7: How to run (Windows)
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.shapes.placeholders[1].text_frame

title.text = "How to Run (Windows PowerShell)"
body.text = "1) Create & activate venv\n2) Install deps\n3) (Optional) Train model\n4) Run Streamlit app"

p = body.add_paragraph()
p.text = "Commands (example):"
p.level = 1

cmds = [
    'python -m venv .venv',
    '.\\.venv\\Scripts\\Activate.ps1',
    'pip install -r requirements.txt',
    'python -m src.train_model',
    'python -m streamlit run src.app.py'
]
for c in cmds:
    p = body.add_paragraph()
    p.text = c
    p.level = 2

# Slide 8: Fixes & Improvements
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.shapes.placeholders[1].text_frame

title.text = "Fixes & Improvements"
body.text = "What was fixed and why:"
fixes = [
    'preprocess_data: use passed DataFrame (avoid reloading)',
    "Mapping yes/no with explicit per-column mapping (removes FutureWarning)",
    'Aligned Streamlit selectbox options with dataset values',
    'Added scaled pipeline and increased max_iter to prevent convergence warnings',
    'Added test script to verify saved model predictions'
]
for f in fixes:
    p = body.add_paragraph()
    p.text = f
    p.level = 1

# Slide 9: Troubleshooting & Notes
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.shapes.placeholders[1].text_frame

title.text = "Troubleshooting & Notes"
notes = [
    'If Python not found: install Python and add to PATH',
    'If package missing: pip install -r requirements.txt',
    'Streamlit runs at http://localhost:8501 by default',
    'Use src/test_model_prediction.py to sanity-check model outputs'
]
body.text = ''
for n in notes:
    p = body.add_paragraph()
    p.text = n
    p.level = 1

# Slide 10: Next Steps
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.shapes.placeholders[1].text_frame

title.text = "Next Steps"
nexts = [
    'Add unit tests for preprocessing and model prediction',
    'Add friendly UI labels and internal mapping to dataset values',
    'Include example inputs and expected outputs for CI',
    'Improve model (feature selection, hyperparameter tuning)'
]
for ns in nexts:
    p = body.add_paragraph()
    p.text = ns
    p.level = 1

# Slide 11: Contact / Repo
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.shapes.placeholders[1].text_frame

title.text = "Repository & Contact"
body.text = "Repository: https://github.com/syedasmarali/student-dropout-predictor\nFile: model/logistic_regression_model.pkl\nPresentation generated at: Student_Dropout_Predictor_Presentation.pptx"

# Save Presentation
def main():
    prs.save('Student_Dropout_Predictor_Presentation.pptx')

if __name__ == '__main__':
    main()
